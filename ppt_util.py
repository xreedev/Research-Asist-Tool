from pptx import Presentation

def create_ppt(sentences):
    # Create a presentation object
    prs = Presentation()

    # Set the maximum characters allowed on a slide
    max_chars_per_slide = 10

    # Create an empty slide
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use index 1 for a blank slide layout
    title = slide.shapes.title
    title.text = "Scientific Paper Summary"

    # Initialize current text
    current_text = ""

    # Iterate over sentences
    for sentence in sentences:
        # Check if adding the current sentence exceeds the maximum characters allowed per slide
        if len(current_text) + len(sentence) <= max_chars_per_slide:
            current_text += sentence + " "
        else:
            # Add the current text to the slide
            content = slide.placeholders[1]
            content.text = current_text

            # Create a new slide
            slide = prs.slides.add_slide(prs.slide_layouts[1])
            title = slide.shapes.title
            title.text = "Scientific Paper Summary"

            # Reset current text with the current sentence
            current_text = sentence + " "

    # Add the remaining text to the last slide
    content = slide.placeholders[1]
    content.text = current_text

    return prs
