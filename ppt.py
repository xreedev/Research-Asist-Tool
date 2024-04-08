from pptx import Presentation
# This function adds the required text to the required row
def setText(shape, text):
    print(len(shape.text_frame.paragraphs),len(shape.text_frame.paragraphs[0].runs))
    shape.text_frame.paragraphs[0].runs[0].text = text[0]
    shape.text_frame.paragraphs[2].runs[0].text = text[1]
    shape.text_frame.paragraphs[4].runs[0].text = text[2]
def modifyIntro(intro,body):  #This fucntion modifies the Introduction slide
    body_shape=intro.shapes[3]
    setText(body_shape,body)

def modifyWMRD(rw,body,):  #This function modifies the Related Work,Modification,Results and Discussion slide
    body_shape=rw.shapes[4]
    setText(body_shape,body)

def modifyConcr(conc,body): #This fucntion modifies the conclusion slide
    body_shape=conc.shapes[1]
    body_shape.text_frame.paragraphs[0].runs[0].text = body[0]
    body_shape.text_frame.paragraphs[2].runs[0].text = body[1]

def setTITAuth(title_slide,title_name,author_list): #This fucntion is used to add authors name and heaing to first slide
    title = title_slide.shapes[15]
    author = title_slide.shapes[16]
    title.text_frame.paragraphs[0].runs[0].text = title_name
    name = ""
    for authors in author_list: #create one string to represent all authors
        name += authors
        name += ","
    author.text_frame.paragraphs[0].runs[0].text = name[:-1] #remove last comma and add it to ppt text frame
def save_ppt(presentation_path,dict,title_name,author_list):
    prs = Presentation(presentation_path)
    title_slide = prs.slides[0] #Represent each slide in ppt
    intro=prs.slides[1]
    rel_work= prs.slides[2]
    methodology= prs.slides[3]
    results = prs.slides[4]
    discussion = prs.slides[5]
    conclusion = prs.slides[6]

    setTITAuth(title_slide,title_name,author_list)

    modifyIntro(intro,dict['Introduction'].split("."))
    modifyWMRD(rel_work,dict['Related Work'].split("."))
    modifyWMRD(methodology,dict['Methodology'].split("."))
    modifyWMRD(results,dict['Results'].split("."))
    modifyWMRD(discussion,dict['Discussion'].split("."))
    modifyConcr(conclusion,dict['Conclusion'].split("."))
    prs.save("new.pptx")