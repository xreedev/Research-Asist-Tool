import PyPDF2

def pdf_to_text(pdf_path):
    text = ""
    with open(pdf_path, "rb") as file:
        # Create a PDF reader object
        pdf_reader = PyPDF2.PdfReader(file)

        # Iterate through each page
        for page_num in range(len(pdf_reader.pages)):
            # Extract text from the page
            page = pdf_reader.pages[page_num]
            text += page.extract_text()

    return text

# Path to the PDF file
pdf_path = "C:\\Users\\dell\\Pictures\\Desktop\\btech\\Main Project\\AI.pdf"

# Extract text from the PDF file
text = pdf_to_text(pdf_path)

# Print the extracted text
print(text)


from pptx import Presentation

def create_ppt(text):
    # Create a presentation object
    prs = Presentation()

    # Set the maximum characters allowed on a slide
    max_chars_per_slide = 500

    # Split the text into slides based on maximum characters
    slides_text = [text[i:i+max_chars_per_slide] for i in range(0, len(text), max_chars_per_slide)]

    # Add each slide to the presentation
    for slide_text in slides_text:
        # Add a new slide
        slide = prs.slides.add_slide(prs.slide_layouts[1])  # Use index 1 for a blank slide layout

        # Add a title to the slide
        title = slide.shapes.title
        title.text = "Scientific Paper Summary"

        # Add text to the slide
        content = slide.placeholders[1]
        content.text = slide_text

    return prs

# Create a PowerPoint presentation from the extracted text
presentation = create_ppt(text)

# Save the presentation to a file
presentation.save("test.pptx")


